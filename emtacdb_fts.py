from plugins.ai_models import generate_embedding, store_embedding
import base64
import json
import os
import re
import time
import logging
import comtypes.client
import win32com.client

from datetime import datetime
from io import BytesIO
from io import BytesIO as BytesIOImage
from pathlib import Path

import docx
from docx2pdf import convert
import fitz
import numpy as np
import openai
import pandas as pd
import pdfplumber
import pythoncom
import requests
import spacy
from fuzzywuzzy import fuzz, process
from nltk.corpus import wordnet
from PIL import Image as PILImage
from pptx import Presentation
from werkzeug.security import check_password_hash, generate_password_hash
from flask import current_app, send_file, url_for
from sqlalchemy.orm import configure_mappers
from sqlalchemy import (DateTime, Date,Column, ForeignKey, Integer, JSON, LargeBinary, Enum as SqlEnum,
                        String, Table, and_, create_engine, func, or_, text, Float, Text, UniqueConstraint)
from enum import Enum as PyEnum  # Import Enum and alias it as PyEnum
from sqlalchemy.exc import SQLAlchemyError
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.ext.mutable import MutableList
from sqlalchemy.orm import declarative_base, configure_mappers, relationship, scoped_session, sessionmaker
from config import (CURRENT_EMBEDDING_MODEL, OPENAI_API_KEY, BASE_DIR, COPY_FILES, DATABASE_DIR, DATABASE_DOC,
                    DATABASE_PATH, DATABASE_URL, KEYWORDS_FILE_PATH, TEMPORARY_FILES,
                    TEMPORARY_UPLOAD_FILES, UPLOAD_FOLDER,DATABASE_PATH_IMAGES_FOLDER)

# Configure mappers (must be called after all ORM classes are defined)
configure_mappers()

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load the English language model
nlp = spacy.load('en_core_web_sm')

# Set your OpenAI API key
openai.api_key = OPENAI_API_KEY

# Constants for chunk size and model name
CHUNK_SIZE = 8000
MODEL_NAME = "text-embedding-ada-002"

# Constants for the database directory and path
DATABASE_DIR = os.path.join(BASE_DIR, 'Database')
DATABASE_PATH = os.path.join(DATABASE_DIR, 'emtac_db.db')

# Check if the database file exists in the specified directory and create it if not
if not os.path.exists(DATABASE_PATH):
    open(DATABASE_PATH, 'w').close()

# Create SQLAlchemy engine
engine = create_engine(DATABASE_URL)
Base = declarative_base()
Session = scoped_session(sessionmaker(bind=engine))  # Use scoped_session here
session = Session()

Base = declarative_base()

# Main Tables
class SiteLocation(Base):
    __tablename__ = 'site_location'
    id = Column(Integer, primary_key=True)
    title = Column(String, nullable=False)
    room_number = Column(String, nullable=False)
    
    position = relationship('Position', back_populates="site_location")

class Position(Base):
    __tablename__ = 'position'
    id = Column(Integer, primary_key=True)
    area_id = Column(Integer, ForeignKey('area.id'), nullable=True)
    equipment_group_id = Column(Integer, ForeignKey('equipment_group.id'), nullable=True)
    model_id = Column(Integer, ForeignKey('model.id'), nullable=True)
    asset_number_id = Column(Integer, ForeignKey('asset_number.id'), nullable=True)
    location_id = Column(Integer, ForeignKey('location.id'), nullable=True)
    site_location_id = Column(Integer, ForeignKey('site_location.id'), nullable=True)

    area = relationship("Area", back_populates="position")
    equipment_group = relationship("EquipmentGroup", back_populates="position")
    model = relationship("Model", back_populates="position")
    asset_number = relationship("AssetNumber", back_populates="position")
    location = relationship("Location", back_populates="position")
    bill_of_material = relationship("BillOfMaterial", back_populates="position")
    part_position = relationship("PartsPositionAssociation", back_populates="position")
    image_position_association = relationship("ImagePositionAssociation", back_populates="position")
    drawing_position = relationship("DrawingPositionAssociation", back_populates="position")
    problem_position = relationship("ProblemPositionAssociation", back_populates="position")
    completed_document_position_association = relationship("CompletedDocumentPositionAssociation", back_populates="position")
    site_location = relationship("SiteLocation", back_populates="position")

     
class Area(Base):
    __tablename__ = 'area'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)
    description = Column(String)
    
    equipment_group = relationship("EquipmentGroup", back_populates="area")
    position = relationship("Position", back_populates="area")
   
class EquipmentGroup(Base):
    __tablename__ = 'equipment_group'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)
    area_id = Column(Integer, ForeignKey('area.id'))
    
    area = relationship("Area", back_populates="equipment_group") 
    model = relationship("Model", back_populates="equipment_group")
    position = relationship("Position", back_populates="equipment_group")
    
class Model(Base):
    __tablename__ = 'model'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)
    description = Column(String)
    equipment_group_id = Column(Integer, ForeignKey('equipment_group.id'))
    
    equipment_group = relationship("EquipmentGroup", back_populates="model")
    asset_number = relationship("AssetNumber", back_populates="model")
    location = relationship("Location", back_populates="model")
    position = relationship("Position", back_populates="model")
    
class AssetNumber(Base):
    __tablename__ = 'asset_number'

    id = Column(Integer, primary_key=True)
    number = Column(String, nullable=False)    
    description = Column(String)
    model_id = Column(Integer, ForeignKey('model.id'))
    
    model = relationship("Model", back_populates="asset_number")
    position = relationship("Position", back_populates="asset_number")
    
class Location(Base):
    __tablename__ = 'location'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)    
    model_id = Column(Integer, ForeignKey('model.id'))
    
    model = relationship("Model", back_populates="location")
    position = relationship("Position", back_populates="location")
   
class Part(Base):
    __tablename__ = 'part'

    id = Column(Integer, primary_key=True)
    part_number = Column(String, unique=True)  # ITEMNUM
    name = Column(String)  # DESCRIPTION
    oem_mfg = Column(String)  # OEMMFG
    model = Column(String)  # MODEL
    class_flag = Column(String)  # Class Flag
    ud6 = Column(String)  # UD6
    type = Column(String)  # TYPE
    notes = Column(String)  # Notes
    documentation = Column(String)  # Specifications

    part_position = relationship("PartsPositionAssociation", back_populates="part")
    bill_of_material = relationship("BillOfMaterial", back_populates="part")
    part_problem = relationship("PartProblemAssociation", back_populates="part")
    part_solution = relationship("PartSolutionAssociation", back_populates="part")
    drawing_part = relationship("DrawingPartAssociation", back_populates="part")

    __table_args__ = (UniqueConstraint('part_number', name='_part_number_uc'),)

class Image(Base):
    __tablename__ = 'image'

    id = Column(Integer, primary_key=True)
    title = Column(String, nullable=False)
    description = Column(String, nullable=False)
    file_path = Column(String, nullable=False)
    
    image_position_association = relationship("ImagePositionAssociation", back_populates="image")
    image_problem = relationship("ImageProblemAssociation", back_populates="image")
    image_solution = relationship("ImageSolutionAssociation", back_populates="image")
    bill_of_material = relationship("BillOfMaterial", back_populates="image")
    image_completed_document_association = relationship("ImageCompletedDocumentAssociation", back_populates="image")
    image_embedding = relationship("ImageEmbedding", back_populates="image")

class ImageEmbedding(Base):
    __tablename__ = 'image_embedding'

    id = Column(Integer, primary_key=True)
    image_id = Column(Integer, ForeignKey('image.id'))
    model_name = Column(String, nullable=False)
    model_embedding = Column(LargeBinary, nullable=False)

    image = relationship("Image", back_populates="image_embedding")

class Drawing(Base):
    __tablename__ = 'drawing'

    id = Column(Integer, primary_key=True)
    drw_equipment_name = Column(String)
    drw_number = Column(String)
    drw_name = Column(String)
    drw_revision = Column(String)
    drw_spare_part_number = Column(String)
    file_path = Column(String, nullable=False)
    
    drawing_position = relationship("DrawingPositionAssociation", back_populates="drawing")
    drawing_problem = relationship("DrawingProblemAssociation", back_populates="drawing")
    drawing_solution = relationship("DrawingSolutionAssociation", back_populates="drawing")
    drawing_part = relationship("DrawingPartAssociation", back_populates="drawing")

class Document(Base):
    __tablename__ = 'document'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)
    file_path = Column(String, nullable=False)
    content = Column(String)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    embedding = Column(LargeBinary)
    rev = Column(String)
    
    embeddings = relationship("DocumentEmbedding", back_populates="document")
    complete_document = relationship("CompleteDocument", back_populates="document")

class DocumentEmbedding(Base):
    __tablename__ = 'document_embedding'

    id = Column(Integer, primary_key=True)
    document_id = Column(Integer, ForeignKey('document.id'))
    model_name = Column(String, nullable=False)
    model_embedding = Column(LargeBinary, nullable=False)

    document = relationship("Document", back_populates="embeddings")

class CompleteDocument(Base):
    __tablename__ = 'complete_document'

    id = Column(Integer, primary_key=True)
    title = Column(String)
    file_path = Column(String)
    content = Column(Text)
    rev = Column(String)

    document = relationship("Document", back_populates="complete_document")
    completed_document_position_association = relationship("CompletedDocumentPositionAssociation", back_populates="complete_document")
    powerpoint = relationship("PowerPoint", back_populates="complete_document")
    image_completed_document_association = relationship("ImageCompletedDocumentAssociation", back_populates="complete_document")
    complete_document_problem = relationship("CompleteDocumentProblemAssociation", back_populates="complete_document")
    complete_document_solution = relationship("CompleteDocumentSolutionAssociation", back_populates="complete_document")

class Problem(Base):
    __tablename__ = 'problem'

    id = Column(Integer, primary_key=True)
    name = Column(String, nullable=False)
    description = Column(String, nullable=False)
    
    solution = relationship("Solution", back_populates="problem")
    problem_position = relationship("ProblemPositionAssociation", back_populates="problem")
    image_problem = relationship("ImageProblemAssociation", back_populates="problem")
    complete_document_problem = relationship("CompleteDocumentProblemAssociation", back_populates="problem")
    drawing_problem = relationship("DrawingProblemAssociation", back_populates="problem")
    part_problem = relationship("PartProblemAssociation", back_populates="problem")

class Solution(Base):
    __tablename__ = 'solution'

    id = Column(Integer, primary_key=True)
    description = Column(String, nullable=False)
    problem_id = Column(Integer, ForeignKey('problem.id'))
    
    problem = relationship("Problem", back_populates="solution")
    image_solution = relationship("ImageSolutionAssociation", back_populates="solution")
    complete_document_solution = relationship("CompleteDocumentSolutionAssociation", back_populates="solution")
    drawing_solution = relationship("DrawingSolutionAssociation", back_populates="solution")
    part_solution = relationship("PartSolutionAssociation", back_populates="solution")

# Class representing PowerPoint presentations in the database
class PowerPoint(Base):
    __tablename__ = 'powerpoint'

    id = Column(Integer, primary_key=True)
    title = Column(String, nullable=False)
    ppt_file_path = Column(String, nullable=False)
    pdf_file_path = Column(String, nullable=False)
    description = Column(String, nullable=True)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    rev = Column(String, nullable=True)

    complete_document = relationship("CompleteDocument", back_populates="powerpoint")

    def __init__(self, title, ppt_file_path, pdf_file_path, complete_document_id, description=None):
        self.title = title
        self.ppt_file_path = ppt_file_path
        self.pdf_file_path = pdf_file_path
        self.complete_document_id = complete_document_id
        self.description = description

# Junction Classes 
class DrawingPartAssociation(Base):
    __tablename__ = 'drawing_part'
    id = Column(Integer, primary_key=True)
    drawing_id = Column(Integer, ForeignKey('drawing.id'))
    part_id = Column(Integer, ForeignKey('part.id'))
    
    drawing = relationship("Drawing", back_populates="drawing_part")
    part = relationship("Part", back_populates="drawing_part")
    
class PartProblemAssociation(Base):
    __tablename__ = 'part_problem'
    id = Column(Integer, primary_key=True)
    part_id = Column(Integer, ForeignKey('part.id'))
    problem_id = Column(Integer, ForeignKey('problem.id'))
    
    part = relationship("Part", back_populates="part_problem")
    problem = relationship("Problem", back_populates="part_problem")

class PartSolutionAssociation(Base):
    __tablename__ = 'part_solution'
    id = Column(Integer, primary_key=True)
    part_id = Column(Integer, ForeignKey('part.id'))
    solution_id = Column(Integer, ForeignKey('solution.id'))
    
    part = relationship("Part", back_populates="part_solution")
    solution = relationship("Solution", back_populates="part_solution")
   
class DrawingProblemAssociation(Base):
    __tablename__ = 'drawing_problem'
    id = Column(Integer, primary_key=True)
    drawing_id = Column(Integer, ForeignKey('drawing.id'))
    problem_id = Column(Integer, ForeignKey('problem.id'))
    
    drawing = relationship("Drawing", back_populates="drawing_problem")
    problem = relationship("Problem", back_populates="drawing_problem")

class DrawingSolutionAssociation(Base):
    __tablename__ = 'drawing_solution'
    id = Column(Integer, primary_key=True)
    drawing_id = Column(Integer, ForeignKey('drawing.id'))
    solution_id = Column(Integer, ForeignKey('solution.id'))
    
    drawing = relationship("Drawing", back_populates="drawing_solution")
    solution = relationship("Solution", back_populates="drawing_solution")
    
class BillOfMaterial(Base):
    __tablename__ = 'bill_of_material'
    id = Column(Integer, primary_key=True)
    part_id = Column(Integer, ForeignKey('part.id'))
    position_id = Column(Integer, ForeignKey('position.id'))
    image_id = Column(Integer, ForeignKey('image.id'))
    quantity = Column(Float, nullable=False)  # Corrected to Float
    comment = Column(String)
    
    part = relationship("Part", back_populates="bill_of_material")
    position = relationship("Position", back_populates="bill_of_material")
    image = relationship("Image", back_populates="bill_of_material")
    
class ProblemPositionAssociation(Base):
    __tablename__ = 'problem_position'
    id = Column(Integer, primary_key=True)
    problem_id = Column(Integer, ForeignKey('problem.id'))
    position_id = Column(Integer, ForeignKey('position.id'))
    
    problem = relationship("Problem", back_populates="problem_position")
    position = relationship("Position", back_populates="problem_position")

class CompleteDocumentProblemAssociation(Base):
    __tablename__ = 'complete_document_problem'
    
    id = Column(Integer, primary_key=True)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    problem_id = Column(Integer, ForeignKey('problem.id'))
    
    complete_document = relationship("CompleteDocument", back_populates="complete_document_problem")
    problem = relationship("Problem", back_populates="complete_document_problem")
    
class CompleteDocumentSolutionAssociation(Base):
    __tablename__ = 'complete_document_solution'
    
    id = Column(Integer, primary_key=True)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    solution_id = Column(Integer, ForeignKey('solution.id'))
    
    complete_document = relationship("CompleteDocument", back_populates="complete_document_solution")
    solution = relationship("Solution", back_populates="complete_document_solution")

class ImageProblemAssociation(Base):
    __tablename__ = 'image_problem'
    id = Column(Integer, primary_key=True)
    image_id = Column(Integer, ForeignKey('image.id'))
    problem_id = Column(Integer, ForeignKey('problem.id'))
    
    image = relationship("Image", back_populates="image_problem")
    problem = relationship("Problem", back_populates="image_problem")

class ImageSolutionAssociation(Base):
    __tablename__ = 'image_solution'
    id = Column(Integer, primary_key=True)
    image_id = Column(Integer, ForeignKey('image.id'))
    solution_id = Column(Integer, ForeignKey('solution.id'))
    
    image = relationship("Image", back_populates="image_solution")
    solution = relationship("Solution", back_populates="image_solution")

class PartsPositionAssociation(Base):
    __tablename__ = 'part_position'
    id = Column(Integer, primary_key=True)
    part_id = Column(Integer, ForeignKey('part.id'))
    position_id = Column(Integer, ForeignKey('position.id'))

    part = relationship("Part", back_populates="part_position")
    position = relationship("Position", back_populates="part_position")
    
class ImagePositionAssociation(Base):
    __tablename__ = 'image_position_association'
    id = Column(Integer, primary_key=True)
    image_id = Column(Integer, ForeignKey('image.id'))
    position_id = Column(Integer, ForeignKey('position.id'))
    
    image = relationship("Image", back_populates="image_position_association")
    position = relationship("Position", back_populates="image_position_association")

class DrawingPositionAssociation(Base):
    __tablename__ = 'drawing_position'
    id = Column(Integer, primary_key=True)
    drawing_id = Column(Integer, ForeignKey('drawing.id'))
    position_id = Column(Integer, ForeignKey('position.id'))
    
    drawing = relationship("Drawing", back_populates="drawing_position")
    position = relationship("Position", back_populates="drawing_position")

class CompletedDocumentPositionAssociation(Base):
    __tablename__ = 'completed_document_position_association'
    id = Column(Integer, primary_key=True)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    position_id = Column(Integer, ForeignKey('position.id'))

    complete_document = relationship("CompleteDocument", back_populates="completed_document_position_association")
    position = relationship("Position", back_populates="completed_document_position_association")

class ImageCompletedDocumentAssociation(Base):
    __tablename__ = 'image_completed_document_association'

    id = Column(Integer, primary_key=True)
    complete_document_id = Column(Integer, ForeignKey('complete_document.id'))
    image_id = Column(Integer, ForeignKey('image.id'))
    
    complete_document = relationship("CompleteDocument", back_populates="image_completed_document_association")
    image = relationship("Image", back_populates="image_completed_document_association")
    
# Process Classes
class FileLog(Base):
    __tablename__ = 'file_logs'
    log_id = Column(Integer, primary_key=True, autoincrement=True)
    session = Column(Integer, nullable=False)
    session_datetime = Column(DateTime, nullable=False)
    file_processed = Column(String)  # Added column for file processed
    total_time = Column(String)

class KeywordAction(Base):
    __tablename__ = 'keyword_actions'

    id = Column(Integer, primary_key=True)
    keyword = Column(String, unique=True)
    action = Column(String)

    # Manually define the query attribute
    query = scoped_session(session).query_property()

    @classmethod
    def find_best_match(cls, user_input, session):
        try:
            # Retrieve all keywords from the database
            all_keywords = [keyword.keyword for keyword in session.query(cls).all()]

            # Use fuzzy string matching to find the best matching keyword
            logger.debug("All keywords: %s", all_keywords)
            matched_keyword, similarity_score = process.extractOne(user_input, all_keywords)
            logger.debug("Matched keyword: %s", matched_keyword)
            logger.debug("Similarity score: %s", similarity_score)

            # Set a threshold for the minimum similarity score
            threshold = 50

            # If the similarity score exceeds the threshold, return the matched keyword and its associated action
            if similarity_score >= threshold:
                # Extract keyword and details using spaCy
                keyword, details = extract_keyword_and_details(user_input)
                if keyword:
                    # Retrieve the associated action from the database using the matched keyword
                    keyword_entry = session.query(cls).filter_by(keyword=keyword).first()
                    if keyword_entry:
                        action = keyword_entry.action
                        logger.debug("Associated action: %s", action)
                        return keyword, action, None  # No need to extract details

            # If no matching keyword is found or similarity score is below threshold, return None
            logger.debug("No matching keyword found or similarity score is below threshold.")
            return None, None, None

        except SQLAlchemyError as e:
            # Handle SQLAlchemy errors
            logger.error("Database error: %s", e)
            return None, None, None

        except Exception as e:
            # Handle other unexpected errors
            logger.error("Unexpected error: %s", e)
            return None, None, None

class ChatSession(Base):
    __tablename__ = 'chat_sessions'
    session_id = Column(Integer, primary_key=True, autoincrement=True)
    user_id = Column(String, nullable=False)
    start_time = Column(String, nullable=False)
    last_interaction = Column(String, nullable=False)
    session_data = Column(MutableList.as_mutable(JSON), default=[])  # Initialize as empty list
    conversation_summary = Column(MutableList.as_mutable(JSON), default=[])  # New column for conversation summary

    def __init__(self, user_id, start_time, last_interaction, session_data=None, conversation_summary=None):
        self.user_id = user_id
        self.start_time = start_time
        self.last_interaction = last_interaction
        if session_data is None:
            session_data = []  # Initialize session_data as empty list if not provided
        self.session_data = session_data
        if conversation_summary is None:
            conversation_summary = []  # Initialize conversation_summary as empty list if not provided
        self.conversation_summary = conversation_summary
      
class QandA(Base):
    __tablename__ = 'qanda'
    id = Column(Integer, primary_key=True, autoincrement=True)
    user_id = Column(String)
    question = Column(String)
    answer = Column(String)
    comment = Column(String)
    rating = Column(String)
    timestamp = Column(String, nullable=False)
    
    def __init__(self, user_id, question, answer, timestamp, rating=None, comment=None):
        self.user_id = user_id
        self.question = question
        self.answer = answer
        self.timestamp = timestamp
        self.rating = rating
        self.comment = comment      

class UserLevel(PyEnum):
    ADMIN = 'admin'
    LEVEL_III = 'level_iii'
    LEVEL_II = 'level_ii'
    STANDARD = 'standard'

class AIModelConfig(Base):
    __tablename__ = 'ai_model_config'

    id = Column(Integer, primary_key=True)
    key = Column(String, unique=True, nullable=False)
    value = Column(String, nullable=False)

# Function to load config from the database
def load_config_from_db():
    session = Session()
    ai_model_config = session.query(AIModelConfig).filter_by(key="CURRENT_AI_MODEL").first()
    embedding_model_config = session.query(AIModelConfig).filter_by(key="CURRENT_EMBEDDING_MODEL").first()
    session.close()

    current_ai_model = ai_model_config.value if ai_model_config else "NoAIModel"
    current_embedding_model = embedding_model_config.value if embedding_model_config else "NoEmbeddingModel"

    return current_ai_model, current_embedding_model

class ImageModelConfig(Base):
    __tablename__ = 'image_model_config'

    id = Column(Integer, primary_key=True)
    key = Column(String, unique=True, nullable=False)
    value = Column(String, nullable=False)

def load_image_model_config_from_db():
    session = Session()
    image_model_config = session.query(ImageModelConfig).filter_by(key="CURRENT_IMAGE_MODEL").first()
    session.close()

    current_image_model = image_model_config.value if image_model_config else "no_model"

    return current_image_model


# Define the User model
class User(Base):
    __tablename__ = 'users'

    id = Column(Integer, primary_key=True)
    employee_id = Column(String, unique=True, nullable=False)
    first_name = Column(String, nullable=False)
    last_name = Column(String, nullable=False)
    current_shift = Column(String, nullable=True)
    primary_area = Column(String, nullable=True)
    age = Column(Integer, nullable=True)
    education_level = Column(String, nullable=True)
    start_date = Column(Date, nullable=True)
    hashed_password = Column(String, nullable=False)
    user_level = Column(SqlEnum(UserLevel), default=UserLevel.STANDARD, nullable=False)

    def set_password(self, password):
        self.hashed_password = generate_password_hash(password)

    def check_password_hash(self, password):
        return check_password_hash(self.hashed_password, password)        
       
# Bind the engine to the Base class
Base.metadata.bind = engine

# Then call create_all()
Base.metadata.create_all(engine, checkfirst=True)

# Create the 'documents_fts' table for full-text search
with Session() as session:
    sql_statement = text("CREATE VIRTUAL TABLE IF NOT EXISTS documents_fts USING FTS5(title, content)")
    session.execute(sql_statement)
    session.commit()

# Ask for API key if it's empty
if not OPENAI_API_KEY:
    OPENAI_API_KEY = input("Enter your OpenAI API key: ")
    with open('config.py', 'w') as config_file:
        config_file.write(f'BASE_DIR = "{BASE_DIR}"\n')
        config_file.write(f'copy_files = {COPY_FILES}\n')
        config_file.write(f'OPENAI_API_KEY = "{OPENAI_API_KEY}"\n')

def create_position(area_id, equipment_group_id, model_id, asset_number_id, location_id, site_location_id):
    with Session() as session:
        try:
            # Retrieve the related entities by their IDs
            area_entity = session.query(Area).filter_by(id=area_id).first() if area_id else None
            equipment_group_entity = session.query(EquipmentGroup).filter_by(id=equipment_group_id).first() if equipment_group_id else None
            model_entity = session.query(Model).filter_by(id=model_id).first() if model_id else None
            asset_number_entity = session.query(AssetNumber).filter_by(id=asset_number_id).first() if asset_number_id else None
            location_entity = session.query(Location).filter_by(id=location_id).first() if location_id else None
            site_location_entity = session.query(SiteLocation).filter_by(id=site_location_id).first() if site_location_id else None

            # Check for an existing Position with the same attributes
            existing_position = session.query(Position).filter_by(
                area=area_entity,
                equipment_group=equipment_group_entity,
                model=model_entity,
                asset_number=asset_number_entity,
                location=location_entity,
                site_location=site_location_entity
            ).first()

            if existing_position:
                logger.info(f"Found existing Position with ID: {existing_position.id}")
                return existing_position.id
            else:
                # Create and add the Position entry
                position = Position(
                    area=area_entity,
                    equipment_group=equipment_group_entity,
                    model=model_entity,
                    asset_number=asset_number_entity,
                    location=location_entity,
                    site_location=site_location_entity
                )
                session.add(position)
                session.commit()  # Commit changes to get the position ID
                logger.info(f"Created new Position with ID: {position.id}")
                return position.id

        except Exception as e:
            logger.error(f"An error occurred in create_position: {e}")
            session.rollback()
            return None

# Function to update keywords in the database from an Excel file
def load_keywords_to_db(session: Session):
    try:
        df = pd.read_excel(KEYWORDS_FILE_PATH)
        logger.info(f"Loaded {len(df)} keywords from {KEYWORDS_FILE_PATH}")
        
        for index, row in df.iterrows():
            keyword = row['keyword']
            action = row['action']
            
            existing_keyword = session.query(KeywordAction).filter_by(keyword=keyword).first()
            
            if existing_keyword:
                # Update the action if the keyword already exists
                existing_keyword.action = action
                logger.info(f"Updated keyword: {keyword}")
            else:
                # Create a new keyword if it doesn't exist
                keyword_action = KeywordAction(keyword=keyword, action=action)
                session.add(keyword_action)
                logger.info(f"Added new keyword: {keyword}")
        
        # Commit the changes
        session.commit()
        logger.info("Keyword actions successfully loaded into the database.")
    
    except Exception as e:
        logger.error(f"Error loading keywords to database: {e}")
        session.rollback()
        raise

# Load keywords into the database from the Excel file
#load_keywords_to_db(session)

# Function to preprocess the input text
def preprocess_text(text):
    # Remove common unnecessary words
    text = text.replace("Can you", "").replace("of", "")
    return text.strip()

# Function to preprocess the input text and extract keyword and details using spaCy
def extract_keyword_and_details(text: str):
    try:
        # Preprocess the input text
        text = preprocess_text(text)
        
        # Tokenize the preprocessed text using spaCy
        doc = nlp(text)
        
        # Initialize variables to store keyword and details
        keyword = ""
        details = ""
        
        # Retrieve all keywords from the database
        all_keywords = [keyword_action.keyword for keyword_action in session.query(KeywordAction).all()]

        # Iterate through the tokens
        for token in doc:
            # Check if the token is a keyword
            if token.text in all_keywords:
                keyword += token.text + " "
            else:
                details += token.text + " "
        
        # Remove trailing whitespace
        keyword = keyword.strip()
        details = details.strip()
        
        return keyword, details
    
    except Exception as e:
        logger.error(f"Error extracting keyword and details: {e}")
        raise

# Function to retrieve PowerPoint presentations based on provided criteria
def get_powerpoints_by_title(session: Session, title=None, description=None):
    try:
        # Create a base query for PowerPoint presentations
        query = session.query(PowerPoint)

        # Add filters based on the provided criteria
        if title:
            query = query.filter(PowerPoint.title == title)
        if description:
            query = query.filter(PowerPoint.description == description)

        # Execute the query and retrieve matching PowerPoint presentations
        powerpoint = query.all()
        logger.info(f"Retrieved {len(powerpoint)} PowerPoint presentations matching the criteria.")
        return powerpoint
    except Exception as e:
        logger.error(f"Error while searching PowerPoint presentations: {str(e)}")
        return []
    try:
        with Session() as session:
            # Create a base query for PowerPoint presentations
            query = session.query(PowerPoint)

            # Add filters based on the provided criteria
            filters = []
            if title:
                filters.append(PowerPoint.title == title)
            if area:
                filters.append(PowerPoint.area == area)
            if equipment_group:
                filters.append(PowerPoint.equipment_group == equipment_group)
            if model:
                filters.append(PowerPoint.model == model)
            if asset_number:
                filters.append(PowerPoint.asset_number == asset_number)
            if description:
                filters.append(PowerPoint.description == description)
            
            # Combine the filters with OR condition
            if filters:
                query = query.filter(or_(*filters))
                
            # Execute the query and retrieve matching PowerPoint presentations
            powerpoint = query.all()
            logger.info(f"Retrieved {len(powerpoint)} PowerPoint presentations matching the criteria.")
            return powerpoint
    except Exception as e:
        logger.error(f"Error while searching PowerPoint presentations: {str(e)}")
        return []
        print(f"Error while searching PowerPoint presentations: {str(e)}")
        return []

# Function to serve a PDF file associated with a PowerPoint presentation
def serve_pdf(powerpoint):
    if powerpoint:
        print(f"Debug: Found {len(powerpoint)} PowerPoint presentations")
        # Assuming you want to display the first matching PowerPoint presentation
        powerpoint = powerpoint[0]

        # Construct the full path to the PDF file
        pdf_full_path = os.path.join(PPT2PDF_PDF_FILES_PROCESS, powerpoint.pdf_file_path)
        print(f"Debug: Full PDF path: {pdf_full_path}")

        if os.path.exists(pdf_full_path):
            print("Debug: PDF file exists. Serving the file.")
            # Serve the PDF file as a response
            return send_file(
                pdf_full_path,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f"{powerpoint.title}.pdf"
            )
        else:
            print("Debug: PDF file not found.")
            return "PDF file not found", 404
    else:
        # Handle the case where no PowerPoint presentations are found
        print("Debug: No PowerPoint presentations found.")
        return render_template('powerpoint_search_results.html')

def add_powerpoint_to_db(session: Session, title, ppt_file_path, pdf_file_path, complete_document_id, description=""):
    try:
        # Create a new PowerPoint entry
        new_powerpoint = PowerPoint(
            title=title,
            ppt_file_path=ppt_file_path,
            pdf_file_path=pdf_file_path,
            description=description,
            complete_document_id=complete_document_id
        )
        session.add(new_powerpoint)
        session.commit()

        logger.info(f"Added PowerPoint: {title}")
        return new_powerpoint.id
    except Exception as e:
        session.rollback()
        logger.error(f"An error occurred in add_powerpoint_to_db: {e}")
        return None

# Function to extract images from a PowerPoint presentation
def extract_images_from_ppt(output_dir, ppt_path, ppt_filename, complete_document_id):
    try:
        prs = Presentation(ppt_path)
        image_paths = []

        ppt_filename_only = os.path.basename(ppt_filename)
        title_prefix = os.path.splitext(ppt_filename_only)[0].replace('_', '')

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # Picture shape
                    image = shape.image
                    image_bytes = image.blob
                    title = f"{title_prefix}_image_{slide_index}_{shape_index}.png"
                    image_path = os.path.join(output_dir, title)

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    image_paths.append(image_path)

                    files = {'image': open(image_path, 'rb')}
                    data = {'complete_document_id': complete_document_id}
                    response = requests.post('http://localhost:5000/images/add_image', files=files, data=data)
                    if response.status_code == 200:
                        print("Image processed successfully")
                    else:
                        print(f"Failed to process image: {response.text}")

        return image_paths
    except Exception as e:
        print(f"An error occurred during image extraction from PowerPoint: {e}")
        return []

# Function to serve an image from the database based on its ID
def serve_image(session: Session, image_id):
    logger.info(f"Attempting to serve image with ID: {image_id}")
    try:
        image = session.query(Image).filter_by(id=image_id).first()
        if image:
            logger.debug(f"Image found: {image.title}, File path: {image.file_path}")
            file_path = os.path.join(DATABASE_DIR, image.file_path)
            if os.path.exists(file_path):
                logger.info(f"Serving file: {file_path}")
                return send_file(file_path, mimetype='image/jpeg', as_attachment=True, download_name=f"{image.title}.jpeg")
            else:
                logger.error(f"File not found: {file_path}")
                return "Image file not found", 404
        else:
            logger.error(f"Image not found with ID: {image_id}")
            return "Image not found", 404
    except Exception as e:
        logger.error(f"An error occurred while serving the image: {e}")
        return "Internal Server Error", 500


def get_total_images_count(description=''):
    with Session() as session:
        query = session.query(func.count(Image.id))
        if description:
            query = query.filter(Image.description.like(f"%{description}%"))
        total_count = query.scalar()
    return total_count

# Example usage:
# total_images = get_total_images_count(description="example")

# Function to add an image to the database
def add_image_to_db(title, file_path, position_id=None, completed_document_position_association_id=None, complete_document_id=None, description=""):
    try:
        logger.debug("Debugging add_image_to_db:")
        logger.debug(f"Title: {title}")
        logger.debug(f"File Path: {file_path}")
        logger.debug(f"Position ID: {position_id}")
        logger.debug(f"Completed Document Position Association ID: {completed_document_position_association_id}")
        logger.debug(f"Complete Document ID: {complete_document_id}")
        logger.debug(f"Description: {description}")

        with Session() as session:
            logger.info(f'Processing image: {title}')

            # Check if an image with the same title and description exists
            existing_image = session.query(Image).filter(and_(Image.title == title, Image.description == description)).first()
            if existing_image is not None and existing_image.file_path == file_path:
                # If an image with the same title, description, and file path exists, do not add the new image
                logger.info(f"Image with same title, description, and file path already exists: {title}")
                new_image = existing_image
            else:
                # Create a new Image object with the provided parameters
                new_image = Image(
                    title=title,
                    description=description,
                    file_path=file_path
                )

                # Add the new image to the session
                session.add(new_image)
                session.commit()  # Commit to get the new image ID

                logger.info(f"Added image: {title}")

            if position_id:
                # Create an entry in the ImagePositionAssociation table
                image_position_association = ImagePositionAssociation(
                    image_id=new_image.id,
                    position_id=position_id
                )
                session.add(image_position_association)
                logger.info(f"Created ImagePositionAssociation with image ID {new_image.id} and position ID {position_id}")

            if complete_document_id:
                # Create an entry in the ImageCompletedDocumentAssociation table
                image_completed_document_association = ImageCompletedDocumentAssociation(
                    image_id=new_image.id,
                    complete_document_id=complete_document_id
                )
                session.add(image_completed_document_association)
                logger.info(f"Created ImageCompletedDocumentAssociation with image ID {new_image.id} and completed document ID {complete_document_id}")

            session.commit()

    except Exception as e:
        logger.error(f"An error occurred in add_image_to_db: {e}")
        logger.error(f"Attempted to process image: {title}")

def split_text_into_chunks(text, max_words=300, pad_token=""):
    logger.info("Starting split_text_into_chunks")
    logger.debug(f"Text length: {len(text)}")
    logger.debug(f"Max words per chunk: {max_words}")
    
    chunks = []
    words = re.findall(r'\S+\s*', text)
    logger.debug(f"Total words found: {len(words)}")
    
    current_chunk = []

    for word in words:
        if word.strip() != pad_token:
            current_chunk.append(word)
        
        if len(current_chunk) >= max_words or word.strip() == "":
            # Pad the current chunk to the specified max_words
            while len(current_chunk) < max_words:
                current_chunk.append(pad_token)
            
            # Add the current chunk to the list of chunks
            chunks.append(" ".join(current_chunk))
            logger.debug(f"Added chunk: {' '.join(current_chunk)}")
            
            # Reset the current chunk
            current_chunk = []
    
    # If there is any remaining content, pad and add it as the last chunk
    if current_chunk:
        while len(current_chunk) < max_words:
            current_chunk.append(pad_token)
        chunks.append(" ".join(current_chunk))
        logger.debug(f"Added last chunk: {' '.join(current_chunk)}")

    logger.info(f"Total chunks created: {len(chunks)}")
    return chunks

# Function to generate an embedding for a given text this has moved to the ai_models.py
#def generate_embedding(document_content):
    logger.info("Starting generate_embedding")
    logger.debug(f"Document content length: {len(document_content)}")

    try:
        response = openai.Embedding.create(
            input=document_content,
            model=MODEL_NAME,
            OPENAI_API_KEY=OPENAI_API_KEY
        )
        embeddings = response.data[0].embedding
        logger.info("Successfully generated embedding")
        return embeddings
    except Exception as e:
        logger.error(f"An error occurred while generating embedding: {e}")
        return None

# Function to extract text from a PDF file using pdfplumber
def extract_text_from_pdf(file_path):
    logger.info("Starting extract_text_from_pdf")
    logger.debug(f"File path: {file_path}")

    try:
        pdf_filepath = os.path.join(DATABASE_DOC, file_path)
        logger.debug(f"Full PDF file path: {pdf_filepath}")
        text = ""
        
        with pdfplumber.open(pdf_filepath) as pdf:
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text += page_text
                logger.debug(f"Extracted text from page {page_num}: {len(page_text) if page_text else 0} characters")
        
        logger.info("Successfully extracted text from PDF")
        return text
    except Exception as e:
        logger.error(f"An error occurred while extracting text from PDF: {e}")
        return None

# Function to extract text from a text (.txt) file
def extract_text_from_txt(txt_path):
    logger.info("Starting extract_text_from_txt")
    logger.debug(f"TXT file path: {txt_path}")

    try:
        with open(txt_path, 'r', encoding='utf-8') as txt_file:
            text = txt_file.read()
            logger.debug(f"Extracted text length: {len(text)} characters")
        logger.info("Successfully extracted text from TXT file")
        return text
    except Exception as e:
        logger.error(f"An error occurred while extracting text from TXT file: {e}")
        return None

def add_docx_to_db(title, docx_path, position_id):
    try:
        # Initialize COM
        pythoncom.CoInitialize()
        # Convert Word document to PDF
        pdf_output_path = docx_path[:-5] + ".pdf"  # Change the extension to .pdf
        convert(docx_path, pdf_output_path)
        
        # Now call add_document_to_db with the generated PDF
        complete_document_id, success = add_document_to_db(title, pdf_output_path, position_id)
        
        if success:
            logger.info(f"Successfully added DOCX document: {title} as PDF with ID: {complete_document_id}")
            return True
        else:
            logger.error(f"Failed to add DOCX document: {title} as PDF")
            return False
    except Exception as e:
        logger.error(f"An error occurred in add_docx_to_db: {e}")
        return False
  
def add_document_to_db(title, file_path, position_id):
    try:
        extracted_text = None
        complete_document_id = None
        completed_document_position_association_id = None
        with Session() as session:
            logger.info(f"Processing file: {file_path}")
            if file_path.endswith(".pdf"):
                logger.info("Extracting text from PDF...")
                extracted_text = extract_text_from_pdf(file_path)
                logger.info("Text extracted from PDF.")
            elif file_path.endswith(".txt"):
                logger.info("Extracting text from TXT...")
                extracted_text = extract_text_from_txt(file_path)
                logger.info("Text extracted from TXT.")
            else:
                logger.error(f"Unsupported file format: {file_path}")
                return None, False

            if extracted_text:
                complete_document = CompleteDocument(
                    title=title,
                    file_path=os.path.relpath(file_path, DATABASE_DIR),
                    content=extracted_text
                )
                session.add(complete_document)
                session.commit()
                complete_document_id = complete_document.id
                logger.info(f"Added complete document: {title}, ID: {complete_document_id}")

                completed_document_position_association = CompletedDocumentPositionAssociation(
                    complete_document_id=complete_document_id,
                    position_id=position_id
                )
                session.add(completed_document_position_association)
                session.commit()
                completed_document_position_association_id = completed_document_position_association.id
                logger.info(f"Added CompletedDocumentPositionAssociation for complete document ID: {complete_document_id}, position ID: {position_id}")

                insert_query_fts = "INSERT INTO documents_fts (title, content) VALUES (:title, :content)"
                session.execute(text(insert_query_fts), {"title": title, "content": extracted_text})
                session.commit()
                logger.info("Added document to the FTS table.")

                text_chunks = split_text_into_chunks(extracted_text)
                for i, chunk in enumerate(text_chunks):
                    padded_chunk = ' '.join(split_text_into_chunks(chunk, pad_token="", max_words=150))
                    document = Document(
                        name=f"{title} - Chunk {i+1}",
                        file_path=os.path.relpath(file_path, DATABASE_DIR),
                        content=padded_chunk,
                        complete_document_id=complete_document_id,
                    )
                    session.add(document)
                    session.commit()
                    logger.info(f"Added chunk {i+1} of document: {title}")

                    if CURRENT_EMBEDDING_MODEL != "NoEmbeddingModel":
                        embeddings = generate_embedding(padded_chunk, CURRENT_EMBEDDING_MODEL)
                        if embeddings is None:
                            logger.warning(f"Failed to generate embedding for chunk {i+1} of document: {title}")
                        else:
                            store_embedding(document.id, embeddings, CURRENT_EMBEDDING_MODEL)
                            logger.info(f"Generated and stored embedding for chunk {i+1} of document: {title}")
                    else:
                        logger.info(f"No embedding generated for chunk {i+1} of document: {title} because no model is selected.")
            else:
                logger.error("No text extracted from the document.")
                return None, False

            if file_path.endswith(".pdf"):
                logger.info("Extracting images from PDF...")
                extract_images_from_pdf(file_path, session, complete_document_id, completed_document_position_association_id)
                logger.info("Images extracted from PDF.")

            logger.info(f"Successfully processed file: {file_path}")
            return complete_document_id, True
    except Exception as e:
        logger.error(f"An error occurred in add_document_to_db: {e}")
        logger.error(f"Attempted Processed file: {file_path}")
        return None, False
  
def add_text_file_to_db(title, txt_file_path, position_id):
    try:
        with open(txt_file_path, 'r', encoding='utf-8') as txt_file:
            file_content = txt_file.read()
        
        # Split the text content into chunks
        text_chunks = split_text_into_chunks(file_content)
        
        with Session() as session:
            # Add the complete document to the database
            complete_document = CompleteDocument(
                title=title,
                file_path=os.path.basename(txt_file_path),
                content=file_content,
                position_id=position_id  # Associate with the position ID
            )
            session.add(complete_document)
            session.commit()
            complete_document_id = complete_document.id
            logger.info(f"Added complete document: {title}, ID: {complete_document_id}")

            # Add chunks to the database
            for i, chunk in enumerate(text_chunks):
                padded_chunk = ' '.join(split_text_into_chunks(chunk, pad_token="", max_words=150))
                embedding = generate_embedding(padded_chunk)
                if embedding is not None:
                    embedding_json = json.dumps(embedding)
                    embedding_bytes = embedding_json.encode('utf-8')
                    document = Document(
                        name=f"{title} - Chunk {i+1}",
                        file_path=os.path.basename(txt_file_path),
                        content=padded_chunk,
                        complete_document_id=complete_document_id,
                        embedding=embedding_bytes
                    )
                    session.add(document)
                    session.commit()
                    logger.info(f"Added chunk {i+1} of document: {title}")

                    # Create DocumentPositionAssociation entry
                    document_position = DocumentPositionAssociation(
                        document_id=document.id,
                        position_id=position_id
                    )
                    session.add(document_position)
                    session.commit()
                    logger.info(f"Added DocumentPositionAssociation for document ID: {document.id}, position ID: {position_id}")

                else:
                    logger.error(f"Failed to add chunk {i+1} of document: {title}")
                    return False
            logger.info("Text file added successfully.")
            return True
    except Exception as e:
        logger.error(f"An error occurred while processing text file: {e}")
        return False

    try:
        with open(txt_file_path, 'r', encoding='utf-8') as txt_file:
            file_content = txt_file.read()
        
        # Split the text content into chunks
        text_chunks = split_text_into_chunks(file_content)
        
        with Session() as session:
            # Add the complete document to the database
            complete_document = CompleteDocument(
                title=title,
                file_path=os.path.basename(txt_file_path),
                content=file_content,
                position_id=position_id  # Associate with the position ID
            )
            session.add(complete_document)
            session.commit()
            complete_document_id = complete_document.id
            logger.info(f"Added complete document: {title}, ID: {complete_document_id}")

            # Add chunks to the database
            for i, chunk in enumerate(text_chunks):
                padded_chunk = ' '.join(split_text_into_chunks(chunk, pad_token="", max_words=150))
                embedding = generate_embedding(padded_chunk)
                if embedding is not None:
                    embedding_json = json.dumps(embedding)
                    embedding_bytes = embedding_json.encode('utf-8')
                    document = Document(
                        name=f"{title} - Chunk {i+1}",
                        file_path=os.path.basename(txt_file_path),
                        content=padded_chunk,
                        complete_document_id=complete_document_id,
                        embedding=embedding_bytes,
                        position_id=position_id  # Associate with the position ID
                    )
                    session.add(document)
                    session.commit()
                    logger.info(f"Added chunk {i+1} of document: {title}")
                else:
                    logger.error(f"Failed to add chunk {i+1} of document: {title}")
                    return False
            logger.info("Text file added successfully.")
            return True
    except Exception as e:
        logger.error(f"An error occurred while processing text file: {e}")
        return False

def add_csv_data_to_db(csv_file, position_id, max_words=300, pad_token=""):
    try:
        # Extract the CSV file name without the extension
        file_name = os.path.splitext(os.path.basename(csv_file))[0]

        # Read CSV data into a DataFrame
        df = pd.read_csv(csv_file)
        
        with Session() as session:
            # Add the complete document to the database
            complete_document = CompleteDocument(
                title=file_name,
                file_path=os.path.basename(csv_file),
                position_id=position_id  # Associate with the position ID
            )
            session.add(complete_document)
            session.commit()
            complete_document_id = complete_document.id
            logger.info(f"Added complete document: {file_name}, ID: {complete_document_id}")

            # Add chunks to the database
            for index, row in df.iterrows():
                # Combine all values in the row into a single string
                row_content = ', '.join(map(str, row))
                
                # Split the row content into chunks with padding
                chunks = split_text_into_chunks(row_content, max_words, pad_token)
                
                for i, chunk in enumerate(chunks):
                    embedding = generate_embedding(chunk)
                    
                    if embedding is not None:
                        embedding_json = json.dumps(embedding)  # Serialize as JSON
                        embedding_bytes = embedding_json.encode('utf-8')  # Convert to bytes
                        document = Document(
                            name=f"{file_name} - Row {index+1} - Chunk {i+1}",
                            file_path=os.path.basename(csv_file),
                            content=chunk,
                            complete_document_id=complete_document_id,
                            embedding=embedding_bytes
                        )
                        session.add(document)
                        session.commit()
                        logger.info(f"Added {file_name} - Row {index+1} - Chunk {i+1} to the database")

                        # Create DocumentPositionAssociation entry
                        document_position = DocumentPositionAssociation(
                            document_id=document.id,
                            position_id=position_id
                        )
                        session.add(document_position)
                        session.commit()
                        logger.info(f"Added DocumentPositionAssociation for document ID: {document.id}, position ID: {position_id}")

                        # Also insert data into the full-text search table using parameterized query
                        insert_query = text("INSERT INTO documents_fts (title, content) VALUES (:title, :content)")
                        session.execute(insert_query, {"title": f"{file_name} - Row {index+1} - Chunk {i+1}", "content": chunk})
                    else:
                        logger.error(f"Failed to add {file_name} - Row {index+1} - Chunk {i+1} to the database")
            logger.info("CSV data added successfully.")
            return True
    except Exception as e:
        logger.error(f"An error occurred while processing CSV data: {e}")
        return False

    try:
        # Extract the CSV file name without the extension
        file_name = os.path.splitext(os.path.basename(csv_file))[0]

        # Read CSV data into a DataFrame
        df = pd.read_csv(csv_file)
        
        with Session() as session:
            # Add the complete document to the database
            complete_document = CompleteDocument(
                title=file_name,
                file_path=os.path.basename(csv_file),
                position_id=position_id  # Associate with the position ID
            )
            session.add(complete_document)
            session.commit()
            complete_document_id = complete_document.id
            logger.info(f"Added complete document: {file_name}, ID: {complete_document_id}")

            # Add chunks to the database
            for index, row in df.iterrows():
                # Combine all values in the row into a single string
                row_content = ', '.join(map(str, row))
                
                # Split the row content into chunks with padding
                chunks = split_text_into_chunks(row_content, max_words, pad_token)
                
                for i, chunk in enumerate(chunks):
                    embedding = generate_embedding(chunk)
                    
                    if embedding is not None:
                        embedding_json = json.dumps(embedding)  # Serialize as JSON
                        embedding_bytes = embedding_json.encode('utf-8')  # Convert to bytes
                        document = Document(
                            name=f"{file_name} - Row {index+1} - Chunk {i+1}",
                            file_path=os.path.basename(csv_file),
                            content=chunk,
                            complete_document_id=complete_document_id,
                            embedding=embedding_bytes,
                            position_id=position_id  # Associate with the position ID
                        )
                        session.add(document)
                        session.commit()
                        logger.info(f"Added {file_name} - Row {index+1} - Chunk {i+1} to the database")

                        # Also insert data into the full-text search table using parameterized query
                        insert_query = text("INSERT INTO documents_fts (title, content) VALUES (:title, :content)")
                        session.execute(insert_query, {"title": f"{file_name} - Row {index+1} - Chunk {i+1}", "content": chunk})
                    else:
                        logger.error(f"Failed to add {file_name} - Row {index+1} - Chunk {i+1} to the database")
            logger.info("CSV data added successfully.")
            return True
    except Exception as e:
        logger.error(f"An error occurred while processing CSV data: {e}")
        return False
       
def cosine_similarity(vector1, vector2):
    dot_product = np.dot(vector1, vector2)
    norm_vector1 = np.linalg.norm(vector1)
    norm_vector2 = np.linalg.norm(vector2)
    return dot_product / (norm_vector1 * norm_vector2)

def search_documents_db(session: Session, title='', area='', equipment_group='', model='', asset_number='', location=''):
    logger.info("Starting search_documents_db")
    logger.debug(f"Search parameters - title: {title}, area: {area}, equipment_group: {equipment_group}, model: {model}, asset_number: {asset_number}, location: {location}")

    try:
        # Use SQLAlchemy to search for complete documents that match the query
        query = session.query(CompleteDocument).join(CompletedDocumentPositionAssociation).join(Position).options(joinedload(CompleteDocument.completed_document_position_association).joinedload(CompletedDocumentPositionAssociation.position))

        # Apply filters based on provided parameters
        if title:
            query = query.filter(CompleteDocument.title.ilike(f'%{title}%'))
        if area:
            query = query.filter(Position.area_id == int(area))
        if equipment_group:
            query = query.filter(Position.equipment_group_id == int(equipment_group))
        if model:
            query = query.filter(Position.model_id == int(model))
        if asset_number:
            query = query.filter(Position.asset_number_id == int(asset_number))
        if location:
            query = query.filter(Position.location_id == int(location))

        results = query.all()

        # Convert the results to a list of dictionaries for JSON response
        documents = [
            {
                'id': doc.id,
                'title': doc.title,
                'content': doc.content,
                'area': doc.completed_document_position_association[0].position.area_id if doc.completed_document_position_association and doc.completed_document_position_association[0].position else None,
                'equipment_group': doc.completed_document_position_association[0].position.equipment_group_id if doc.completed_document_position_association and doc.completed_document_position_association[0].position else None,
                'model': doc.completed_document_position_association[0].position.model_id if doc.completed_document_position_association and doc.completed_document_position_association[0].position else None,
                'asset_number': doc.completed_document_position_association[0].position.asset_number_id if doc.completed_document_position_association and doc.completed_document_position_association[0].position else None,
                'location': doc.completed_document_position_association[0].position.location_id if doc.completed_document_position_association and doc.completed_document_position_association[0].position else None
            }
            for doc in results
        ]

        logger.info(f"Found {len(documents)} documents matching the criteria")
        return {"documents": documents}
    except Exception as e:
        logger.error(f"An error occurred while searching documents: {e}")
        return {"error": str(e)}

        logger.error(f"An error occurred while searching documents: {e}")
        return {"error": str(e)}

# Function to find the most relevant document based on a question
def find_most_relevant_document(question, session: Session):
    if CURRENT_EMBEDDING_MODEL == "NoEmbeddingModel":
        logger.info("Embeddings are disabled. Returning None for document search.")
        return None

    try:
        # Generate embedding for the question
        question_embedding = generate_embedding(question, CURRENT_EMBEDDING_MODEL)
        if not question_embedding:
            logger.info("No embeddings generated. Returning None.")
            return None

        # Fetch documents with the current embedding model
        documents = session.query(Document).join(DocumentEmbedding).filter(DocumentEmbedding.model_name == CURRENT_EMBEDDING_MODEL).all()

        most_relevant_document = None
        highest_similarity = -1

        for doc in documents:
            for embedding_record in doc.embeddings:
                if embedding_record.model_name == CURRENT_EMBEDDING_MODEL:
                    doc_embedding = json.loads(embedding_record.model_embedding.decode('utf-8'))
                    similarity = cosine_similarity(question_embedding, doc_embedding)

                    logger.debug(f"Similarity for document {doc.id} with model {CURRENT_EMBEDDING_MODEL}: {similarity}")

                    if similarity > highest_similarity:
                        highest_similarity = similarity
                        most_relevant_document = doc
                        logger.debug(f"Most relevant document content: {most_relevant_document.content}")

        threshold = 0.01
        if highest_similarity >= threshold:
            logger.info(f"Found most relevant document with ID {most_relevant_document.id} and similarity {highest_similarity}")
            return most_relevant_document
        else:
            logger.info("No relevant document found with sufficient similarity")
            return None
    except Exception as e:
        logger.error(f"An error occurred while finding the most relevant document: {e}")
        return None

def create_session(user_id, session_data, session):
    now = datetime.now().isoformat()
    # Ensure session_data is always a list
    if not isinstance(session_data, list):
        session_data = [session_data]
    new_session = ChatSession(user_id=user_id, start_time=now, last_interaction=now, session_data=[], conversation_summary=[])
    session.add(new_session)
    session.commit()
    return new_session.session_id

def update_session(session_id, session_data, answer, session):
    try:
        now = datetime.now().isoformat()
        print(f"Updating session with ID: {session_id}")
        print("Session Data:", session_data)  # Add this line for debug print
        
        # Ensure session_data is always a list
        if not isinstance(session_data, list):
            session_data = [session_data]
        
        session_to_update = session.query(ChatSession).filter(ChatSession.session_id == session_id).first()
        if session_to_update:
            print("Session found. Updating session...")
            session_to_update.last_interaction = now
            # Append the new question-answer pair to session_data
            session_data_entry = {'question': session_data[-1], 'answer': answer}  # Assuming answer is the bot's response
            session_to_update.session_data.append(session_data_entry)
            
            # Limit session_data to the last 10 entries
            session_to_update.session_data = session_to_update.session_data[-10:]
            
            session.commit()
            print("Session updated successfully.")
    except SQLAlchemyError as e:
        session.rollback()
        logger.error(f"Database error: {e}")
        print(f"Database error occurred: {e}")
    except Exception as e:
        logger.error(f"Unexpected error: {e}")
        print(f"Unexpected error occurred: {e}")

def get_session(user_id, session):
    # Use the passed 'session' for database operations
    return session.query(ChatSession).filter_by(user_id=user_id).first() 

def load_keywords_and_patterns(session: Session):
    keywords_and_patterns = {}
    keyword_to_action_mapping = {}  # Create a dictionary to map keywords to actions

    with current_app.app_context():
        keyword_actions = session.query(KeywordAction).all()
        for ka in keyword_actions:
            # Convert the keyword to lowercase
            keyword = ka.keyword.lower()
            
            # Assuming you can derive a regex pattern from the keyword or action
            # This is a simplified example; you may need to adjust it based on your actual data structure
            pattern = keyword + r" of (.+)"  # Customize this pattern based on your needs
            keywords_and_patterns[keyword] = pattern

            # Associate a single action with each keyword
            keyword_to_action_mapping[keyword] = ka.action

            # Add debug statements to check loaded keywords and patterns
            logger.debug(f"Loaded Keyword: {keyword}")
            logger.debug(f"Pattern: {pattern}")
            logger.debug(f"Action: {ka.action}")

    return keywords_and_patterns, keyword_to_action_mapping

# Function to find a keyword and extract details from a question
def find_keyword_and_extract_detail(question, keywords_and_patterns, keyword_to_action_mapping, session):
    logger.info("Checking if any of the loaded keywords match the user input")

    # Tokenize the input question using spaCy
    doc = nlp(question)

    # Initialize variables to store keyword, details, and action
    matched_keyword = None
    details = None

    # Iterate over the tokens in the input question
    for i, token in enumerate(doc):
        # Lemmatize the token to its base form
        lemma = token.lemma_

        # Check if the lemma is in the keyword mapping
        if lemma in keyword_to_action_mapping:
            matched_keyword = lemma
            # Extract the details by finding the subsequent words
            details = " ".join([t.text for t in doc[i + 1:]])
            logger.debug(f"Matched keyword: {matched_keyword}")
            logger.debug(f"Extracted details: {details}")
            break

    # If no keyword is detected using lemma matching, fall back to traditional keyword matching
    if not matched_keyword:
        # Split the input question into individual words
        words = question.lower().split()
        
        # Iterate over the words in the input question
        for i, word in enumerate(words):
            # Check if the word is a keyword
            if word in keyword_to_action_mapping:
                matched_keyword = word
                # Extract the details by finding the subsequent words
                details = " ".join(words[i + 1:])
                logger.debug(f"Matched keyword: {matched_keyword}")
                logger.debug(f"Extracted details: {details}")
                break

            # Check for pairs of words
            if i < len(words) - 1:
                pair = f"{word} {words[i + 1]}"
                if pair in keyword_to_action_mapping:
                    matched_keyword = pair
                    # Extract the details by finding the subsequent words
                    details = " ".join(words[i + 2:])
                    logger.debug(f"Matched keyword: {matched_keyword}")
                    logger.debug(f"Extracted details: {details}")
                    break

            # Check for triplets of words
            if i < len(words) - 2:
                triplet = f"{word} {words[i + 1]} {words[i + 2]}"
                if triplet in keyword_to_action_mapping:
                    matched_keyword = triplet
                    # Extract the details by finding the subsequent words
                    details = " ".join(words[i + 3:])
                    logger.debug(f"Matched keyword: {matched_keyword}")
                    logger.debug(f"Extracted details: {details}")
                    break

    if matched_keyword:
        logger.info(f"Matched keyword: {matched_keyword}, details: {details}")
        return matched_keyword, details

    try:
        # Request interpretation from OpenAI
        response = openai.Completion.create(
            engine="gpt-3.5-turbo-instruct",
            prompt=f"Is the following question asking to see a document, an image, or a PowerPoint presentation? Additionally, does the question mention any specific title or name? Question: '{question}'",
            temperature=0,
            max_tokens=500,
            top_p=1.0,
            frequency_penalty=0.0,
            presence_penalty=0.0
        )
        interpretation = response.choices[0].text.strip()
        logger.debug(f"OpenAI interpretation: {interpretation}")

        # Extract title from the interpretation
        title = extract_title_from_openai_response(interpretation)
        logger.debug(f"Extracted title: {title}")

        if title is None:
            logger.info("No title extracted from the OpenAI interpretation. Exiting function.")
            return None, None

        # Remove "." and "," from the title
        title = title.replace(".", "").replace(",", "")
        logger.debug(f"Title after processing: {title}")

        # Interpret OpenAI's response and take appropriate action based on the matched keyword
        if "document" in interpretation.lower():
            matched_keyword = "open file"
            details = title
        elif "image" in interpretation.lower():
            matched_keyword = "present photo"
            details = title
        elif "powerpoint" in interpretation.lower():
            matched_keyword = "slide show"
            details = title

        logger.debug(f"Matched keyword from OpenAI interpretation: {matched_keyword}")
        logger.debug(f"Extracted details from OpenAI interpretation: {details}")

        if matched_keyword:
            logger.info(f"Action for keyword: {matched_keyword}, details: {details}")
            return matched_keyword, details
        else:
            logger.info("No action found based on OpenAI interpretation, proceeding with AI logic")
            return None, None

    except Exception as e:
        logger.error(f"Error querying OpenAI: {e}")
        return None, None

def extract_title_from_openai_response(interpretation):
    logger.debug(f"Extracting title from OpenAI interpretation: {interpretation}")
    try:
        logger.debug("Using regular expressions to extract the title from the interpretation")
        title_match = re.search(r'"(.*?)"', interpretation)
        if title_match:
            title = title_match.group(1)
            logger.debug(f"Extracted title before cleanup: {title}")
            # Remove "." and "," from the title
            title = title.replace(".", "").replace(",", "")
            logger.debug(f"Cleaned up title: {title}")
            return title
        else:
            logger.debug("No title match found in the interpretation")
            return None
    except Exception as e:
        logger.error(f"Error occurred while extracting title: {e}")
        return None
        
# Function to perform action based on the matched keyword
def perform_action_based_on_keyword(action, details):
    logger.info("Starting action based on keyword")
    logger.info(f"Action: {action}")
    logger.info(f"Details: {details}")

    if action == "search_images_bp":
        logger.info("Calling the function to show images with the extracted details")
        result = search_images_by_keyword(details)
        logger.debug(f"Show Images Result: {result}")
        return result
    elif action == "search_powerpoints_bp":
        logger.info("Calling the function to get PowerPoint on the specified topic")
        result = search_powerpoints_fts(details)
        logger.debug(f"Get PowerPoint Result: {result}")
        return result
    elif action == "search_documents_bp":
        logger.info("Calling the function to search for documents based on the specified details")
        result = search_documents_fts(details)
        logger.debug(f"Search Documents Result: {result}")
        return result
    else:
        logger.warning("Action not recognized.")
        return "Action not recognized."

    logger.info("End of action based on keyword")

def search_powerpoints_fts(query, session=None):
    try:
        print(f'# Create a SQLAlchemy session if not provided')
        if session is None:
            session = Session()
        
        # Fetch PowerPoint objects from the database based on the search query
        powerpoint = session.query(PowerPoint).filter(PowerPoint.title.ilike(f"%{query}%")).all()

        if powerpoint:
            # Generate HTML anchor tags for each PowerPoint
            html_links = []
            for powerpoint in powerpoint:
                # Generate the relative URL using url_for
                relative_url = url_for('search_powerpoint_fts_bp.view_powerpoint', powerpoint_id=powerpoint.id)

                # Append the base URL to the relative URL
                base_url = 'http://127.0.0.1:5000'
                powerpoint.link = base_url + relative_url

                # Construct HTML anchor tag
                html_links.append(f"Here are your PowerPoint results<a href='{powerpoint.link}'>{powerpoint.title}</a>")

            # Format the search results as a string with HTML anchor tags
            search_results = '\n'.join(html_links)
            return search_results
        else:
            return "No PowerPoint presentations found for your query."
    except Exception as e:
        print(f"Error: {e}")
        return "An error occurred during the search."
    finally:
        # Close the session in the finally block if it's not None
        if session is not None:
            session.close()

def search_documents_fts(query):
    try:
        # Create a SQLAlchemy session
        session = Session()

        # Construct the full-text search query in documents_fts table for the title
        title_search_query = text(
            "SELECT title FROM documents_fts"
        )

        # Execute the query to fetch all titles
        all_titles_results = session.execute(title_search_query)
        all_titles = [row.title for row in all_titles_results]

        # Find matches using fuzzy matching
        matches = []
        for title in all_titles:
            similarity_ratio = fuzz.partial_ratio(query.lower(), title.lower())
            if similarity_ratio >= 80:  # Adjust threshold as needed
                matches.append((title, similarity_ratio))

        # Fetch file paths from complete_document table based on matched titles
        documents = []
        for match in matches:
            title = match[0]
            # Retrieve the corresponding document from the complete_document table
            document = session.query(CompleteDocument).filter_by(title=title).first()
            if document:
                # Generate the relative URL using url_for
                relative_url = url_for('search_documents_fts_bp.view_document', document_id=document.id)

                # Append the base URL to the relative URL
                base_url = 'http://127.0.0.1:5000'
                document.link = base_url + relative_url
                documents.append(document)

        # Close the session
        session.close()

        if documents:
            print(f"Debug: Found {len(documents)} documents")
            # Generate HTML anchor tags for each document
            html_links = generate_html_links(documents)
            # Format the search results as a string with HTML anchor tags
            search_results = '\n'.join(html_links)
            return search_results
        else:
            print("Debug: No documents found.")
            return "No documents found"
    except Exception as e:
        print(f"Error: {e}")
        return "An error occurred during the search."

def search_images_by_keyword(keyword, limit=10, offset=0, session=None):
    try:
        if session is None:
            session = Session()
        print(f'in search_images_by_keyword')
        images = session.query(Image).filter(Image.title.ilike(f'%{keyword}%')).offset(offset).limit(limit).all()
        
        image_data = []
        for image in images:
            thumbnail_src = ""
            try:
                thumbnail = create_thumbnail(os.path.join(DATABASE_DIR,image.file_path))
                if thumbnail:
                    thumbnail_bytes_io = BytesIO()
                    thumbnail.save(thumbnail_bytes_io, format='JPEG')
                    thumbnail_src = f"data:image/jpeg;base64,{base64.b64encode(thumbnail_bytes_io.getvalue()).decode()}"
            except Exception as e:
                logger.error(f"Error creating thumbnail for image ID {image.id}: {e}")

            image_info = {
                'id': image.id,
                'title': image.title,
                'src': f'/serve_image/{image.id}',
                'thumbnail_src': thumbnail_src
            }
            image_data.append(image_info)

        return image_data

    except SQLAlchemyError as e:
        logger.error(f"An error occurred while retrieving images by keyword: {e}")
        return []

    finally:
        if session is not None:
            session.close()

def convert_to_base64(image_blob):
    print(f'# Convert binary image data to Base64 encoded string')
    return base64.b64encode(image_blob).decode('utf-8')

def create_thumbnail(image):
    try:
        # Convert image to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        thumbnail_size = (128, 128)
        image.thumbnail(thumbnail_size)
        return image
    except Exception as e:
        logger.error(f"Error creating thumbnail: {e}")
        return None



def image_to_base64(image):
    print(f'# Convert image to base64 string')
    buffered = BytesIO()
    image.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode("utf-8")   
   
def generate_html_links(documents):
    html_links = []
    for document in documents:
        # Create HTML anchor tag with the document title as the link text and the document link as the href attribute
        html_link = f"Here are you search results <a href='{document.link}'>{document.title}</a>"
        html_links.append(html_link)
    return html_links

#funtion to extract images form a PDF
# Function to extract images from a PDF
def extract_images_from_pdf(file_path, session, complete_document_id, completed_document_position_association_id):
    logger.info(f"Opening PDF file from: {file_path}")
    doc = fitz.open(file_path)
    total_pages = len(doc)
    logger.info(f"Total pages in the PDF: {total_pages}")
    logger.info(f"Inside extract_images_from_pdf, complete_document_id: {complete_document_id}")
    logger.info(f"CompletedDocumentPositionAssociation ID: {completed_document_position_association_id}")
    extracted_images = []

    # Extract file name without extension and remove underscores
    file_name = os.path.splitext(os.path.basename(file_path))[0].replace("_", " ")

    # Ensure the directory exists for temporary upload files
    if not os.path.exists(TEMPORARY_UPLOAD_FILES):
        os.makedirs(TEMPORARY_UPLOAD_FILES)

    for page_num in range(total_pages):
        page = doc[page_num]
        img_list = page.get_images(full=True)

        logger.info(f"Processing page {page_num + 1}/{total_pages} with {len(img_list)} images.")
        
        for img_index, img in enumerate(img_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            # Create a temporary file to send to the Flask route
            temp_path = os.path.join(TEMPORARY_UPLOAD_FILES, f"{file_name}_page{page_num + 1}_image{img_index + 1}.jpg")  # Adjust the extension if needed
            with open(temp_path, 'wb') as temp_file:
                temp_file.write(image_bytes)
            # Log debug information
            logger.info("Sending POST request to http://localhost:5000/image/add_image'")
            logger.info(f"File path: {temp_path}")
            logger.info(f"complete_document_id: {complete_document_id}")
            logger.info(f"completed_document_position_association_id: {completed_document_position_association_id}")
            # Make a POST request to the Flask route to add the image
            response = requests.post('http://localhost:5000/image/add_image', 
                         files={'image': open(temp_path, 'rb')}, 
                         data={'complete_document_id': complete_document_id,
                               'completed_document_position_association_id': completed_document_position_association_id})

            # Check if the request was successful
            if response.status_code == 200:
                logger.info("Image processed successfully")
                # Optionally, you can handle the response if needed
            else:
                logger.error(f"Failed to process image: {response.text}")

            # Wait for a short period before processing the next image
            time.sleep(1)  # Adjust the delay as needed

            # Optionally, you can store the extracted image paths for further processing
            extracted_images.append(temp_path)

    return extracted_images

# Function to get Area object by name
def get_area_by_name(session, area_name):
    try:
        logger.info(f"Fetching Area with name: {area_name}")
        area = session.query(Area).filter_by(name=area_name).first()
        return area
    except Exception as e:
        logger.error(f"Error fetching Area by name: {e}")
        return None

# Function to get EquipmentGroup object by name
def get_equipment_group_by_name(session, equipment_group_name):
    try:
        logger.info(f"Fetching EquipmentGroup with name: {equipment_group_name}")
        equipment_group = session.query(EquipmentGroup).filter_by(name=equipment_group_name).first()
        return equipment_group
    except Exception as e:
        logger.error(f"Error fetching EquipmentGroup by name: {e}")
        return None

# Function to get Model object by name
def get_model_by_name(session, model_name):
    try:
        logger.info(f"Fetching Model with name: {model_name}")
        model = session.query(Model).filter_by(name=model_name).first()
        return model
    except Exception as e:
        logger.error(f"Error fetching Model by name: {e}")
        return None

# Function to get AssetNumber object by number
def get_asset_number_by_number(session, asset_number_number):
    try:
        logger.info(f"Fetching AssetNumber with number: {asset_number_number}")
        asset_number = session.query(AssetNumber).filter_by(number=asset_number_number).first()
        return asset_number
    except Exception as e:
        logger.error(f"Error fetching AssetNumber by number: {e}")
        return None
    
# Function to get Location object by name
def get_location_by_name(session, location_name):
    try:
        logger.info(f"Fetching Location with name: {location_name}")
        location = session.query(Location).filter_by(name=location_name).first()
        return location
    except Exception as e:
        logger.error(f"Error fetching Location by name: {e}")
        return None

# Function to get Area object by ID
def get_area_by_id(session, area_id):
    try:
        logger.info(f"Fetching Area with ID: {area_id}")
        area = session.query(Area).filter_by(id=area_id).first()
        return area
    except Exception as e:
        logger.error(f"Error fetching Area by ID: {e}")
        return None

# Function to get EquipmentGroup object by ID
def get_equipment_group_by_id(session, equipment_group_id):
    try:
        logger.info(f"Fetching EquipmentGroup with ID: {equipment_group_id}")
        equipment_group = session.query(EquipmentGroup).filter_by(id=equipment_group_id).first()
        return equipment_group
    except Exception as e:
        logger.error(f"Error fetching EquipmentGroup by ID: {e}")
        return None

# Function to get Model object by ID
def get_model_by_id(session, model_id):
    try:
        logger.info(f"Fetching Model with ID: {model_id}")
        model = session.query(Model).filter_by(id=model_id).first()
        return model
    except Exception as e:
        logger.error(f"Error fetching Model by ID: {e}")
        return None

# Function to get AssetNumber object by ID
def get_asset_number_by_id(session, asset_number_id):
    try:
        logger.info(f"Fetching AssetNumber with ID: {asset_number_id}")
        asset_number = session.query(AssetNumber).filter_by(id=asset_number_id).first()
        return asset_number
    except Exception as e:
        logger.error(f"Error fetching AssetNumber by ID: {e}")
        return None

# Function to get Location object by ID
def get_location_by_id(session, location_id):
    try:
        logger.info(f"Fetching Location with ID: {location_id}")
        location = session.query(Location).filter_by(id=location_id).first()
        return location
    except Exception as e:
        logger.error(f"Error fetching Location by ID: {e}")
        return None

# Function to check and create directories
def create_directories(directories):
    for directory in directories:
        if not os.path.exists(directory):
            os.makedirs(directory)
            logging.info(f"Created directory: {directory}")
        else:
            logging.info(f"Directory already exists: {directory}")

